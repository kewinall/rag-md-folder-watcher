from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from .base import BaseConverter
from .ocr import UnsupportedImageForOcr, ocr_available, ocr_image_bytes
from ..utils import table_to_markdown


class PptxConverter(BaseConverter):
    extensions = {".pptx"}
    name = "pptx-with-image-ocr"

    def convert(self, path: Path) -> str:
        presentation = Presentation(str(path))
        lines: list[str] = [f"# {path.stem}", ""]
        unsupported_count = 0
        failed_count = 0
        empty_count = 0
        recognized_count = 0
        warning_details: list[str] = []

        for slide_index, slide in enumerate(presentation.slides, 1):
            lines.append(f"## 投影片 {slide_index}")
            picture_index = 0

            for shape in slide.shapes:
                if getattr(shape, "has_text_frame", False):
                    for paragraph in shape.text_frame.paragraphs:
                        text = "".join(run.text for run in paragraph.runs).strip()
                        if text:
                            prefix = "- " if paragraph.level else ""
                            lines.append(f"{prefix}{text}")

                if getattr(shape, "has_table", False):
                    rows = [
                        [cell.text.strip() for cell in row.cells]
                        for row in shape.table.rows
                    ]
                    table_markdown = table_to_markdown(rows)
                    if table_markdown:
                        lines.append(table_markdown)

                if shape.shape_type != MSO_SHAPE_TYPE.PICTURE:
                    continue

                picture_index += 1
                if not (self.config.ocr and ocr_available()):
                    continue

                image = shape.image
                source_name = str(getattr(image, "filename", "") or "")
                content_type = str(getattr(image, "content_type", "") or "")
                try:
                    text = ocr_image_bytes(
                        image.blob,
                        self.config,
                        source_name=source_name,
                        content_type=content_type,
                    )
                except UnsupportedImageForOcr as exc:
                    unsupported_count += 1
                    if len(warning_details) < 10:
                        warning_details.append(
                            f"投影片 {slide_index} 圖片 {picture_index}（{source_name or '未知'}）：{exc}"
                        )
                    continue
                except Exception as exc:
                    failed_count += 1
                    if len(warning_details) < 10:
                        warning_details.append(
                            f"投影片 {slide_index} 圖片 {picture_index}："
                            f"{type(exc).__name__}: {exc}"
                        )
                    continue

                if text:
                    recognized_count += 1
                    lines.extend([
                        f"### 圖片 {picture_index} OCR",
                        text,
                    ])
                else:
                    empty_count += 1

            lines.append("")

        if recognized_count or unsupported_count or failed_count or empty_count:
            self.warnings.append(
                "PowerPoint 圖片 OCR 摘要："
                f"辨識到文字 {recognized_count} 張、無文字 {empty_count} 張、"
                f"略過不支援 {unsupported_count} 張、其他失敗 {failed_count} 張"
            )
        if warning_details:
            self.warnings.append(
                "PowerPoint 圖片 OCR 明細（最多 10 筆）：" + " | ".join(warning_details)
            )

        return "\n".join(lines) + "\n"
